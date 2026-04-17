"""
PowerPoint export for consulting-style Data Storytelling presentations.
Implements the Pyramid Principle and insight compression for McKinsey-style slide generation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, List, Any
import io
import base64

class PPTXStorytellingEngine:
    """
    McKinsey-style PPT Generation Engine.
    Converts data insights into a cohesive narrative structure.
    """

    COLORS = {
        'dark_blue': RGBColor(0x00, 0x2B, 0x5C),     # Primary headers
        'light_blue': RGBColor(0x00, 0x7D, 0xB8),    # Accents
        'gray': RGBColor(0x59, 0x59, 0x59),          # Subtext
        'light_gray': RGBColor(0xF2, 0xF2, 0xF2),    # Background elements
        'black': RGBColor(0x1A, 0x1A, 0x1A),         # Primary text
        'white': RGBColor(0xFF, 0xFF, 0xFF)
    }

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def generate_report(self, report_data: Dict[str, Any], charts: Dict[str, Any] = None) -> bytes:
        """Executes the Storytelling pipeline."""
        
        # 1. Compress Insights (Max 5)
        raw_insights = report_data.get('top_insights', [])
        compressed_insights = sorted(
            raw_insights, 
            key=lambda x: x.get('statistics', {}).get('business_score', 0), 
            reverse=True
        )[:5]

        # 2. Answer First: Title & Executive Summary
        self._add_title_slide(compressed_insights)
        self._add_executive_summary(report_data.get('executive_summary', ''), len(compressed_insights))

        # 3. Supporting Deep-Dive Logic: One slide per insight
        for i, insight in enumerate(compressed_insights, 1):
            self._add_insight_deep_dive(insight, i, charts)

        # 4. Next Steps: Recommendations
        recs = report_data.get('recommendations', [])
        if recs:
            self._add_recommendations_slide(recs)

        # 5. Save Context
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output.getvalue()

    def _add_title_slide(self, top_insights: List[Dict[str, Any]]):
        """Creates an impactful, bold title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Bold Blue Background Strip
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2), self.prs.slide_width, Inches(3.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLORS['dark_blue']
        bg.line.fill.background()

        # Title
        lead_msg = "STRATEGIC DATA REVIEW"
        if top_insights:
            lead_msg = top_insights[0].get('title', lead_msg).upper()
            
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = lead_msg
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.LEFT

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(0.8))
        p = sub_box.text_frame.paragraphs[0]
        p.text = "AI-Generated Strategy & Data Narrative"
        p.font.size = Pt(20)
        p.font.color.rgb = self.COLORS['light_blue']
        p.alignment = PP_ALIGN.LEFT

    def _add_executive_summary(self, summary: str, total_insights: int):
        """Pyramid Principle Level 1: Executive Summary Answer First."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_mckinsey_header(slide, "Executive Summary", "High-level diagnostic and core takeaways")

        # Main text
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(7.5), Inches(4.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = summary if summary else "No significant findings identified."
        p.font.size = Pt(18)
        p.font.color.rgb = self.COLORS['black']
        p.line_spacing = 1.4

        # Side bar (Metric Box)
        stats_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), Inches(1.8), Inches(3.333), Inches(4))
        stats_box.fill.solid()
        stats_box.fill.fore_color.rgb = self.COLORS['light_gray']
        stats_box.line.fill.background()

        p = stats_box.text_frame.paragraphs[0]
        p.text = "\nKey Highlights"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['dark_blue']
        p.alignment = PP_ALIGN.CENTER
        
        p2 = stats_box.text_frame.add_paragraph()
        p2.text = f"\n{total_insights} Core Insights"
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = self.COLORS['black']
        p2.alignment = PP_ALIGN.CENTER

    def _add_insight_deep_dive(self, insight: Dict[str, Any], number: int, charts: Dict[str, Any]):
        """Pyramid Principle Level 2: Supporting arguments / deep dive."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Action-oriented headline
        headline = insight.get('title', f'Insight {number}')
        sub_text = insight.get('narrative', '')
        self._add_mckinsey_header(slide, headline, sub_text)

        # 70% Chart Space
        has_chart = False
        if charts and f'insight_{number}' in charts:
            try:
                img_data = base64.b64decode(charts[f'insight_{number}'])
                slide.shapes.add_picture(io.BytesIO(img_data), Inches(0.5), Inches(2.2), width=Inches(8))
                has_chart = True
            except:
                pass
                
        if not has_chart:
            # Fallback chart box
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(8), Inches(4.5))
            box.fill.solid()
            box.fill.fore_color.rgb = self.COLORS['light_gray']
            box.line.color.rgb = self.COLORS['gray']
            p = box.text_frame.paragraphs[0]
            p.text = "Data Visualization Rendered in Web UI"
            p.font.size = Pt(14)
            p.font.color.rgb = self.COLORS['gray']
            p.alignment = PP_ALIGN.CENTER

        # 30% Callout Space
        why = insight.get('why_it_matters', '')
        if why:
            callout = slide.shapes.add_textbox(Inches(9), Inches(2.2), Inches(3.8), Inches(4.5))
            tf = callout.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = "Why it matters:"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['dark_blue']
            
            p2 = tf.add_paragraph()
            p2.text = f"\n{why}"
            p2.font.size = Pt(14)
            p2.font.color.rgb = self.COLORS['black']
            p2.line_spacing = 1.3

    def _add_recommendations_slide(self, recommendations: List[str]):
        """Pyramid Principle Level 3: Actionable Next Steps."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_mckinsey_header(slide, "Strategic Action Plan", "Recommended next steps to drive impact")

        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
        tf = text_box.text_frame
        tf.word_wrap = True

        for i, rec in enumerate(recommendations[:5], 1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            p.text = f"Priority {i}:  {rec}"
            p.font.size = Pt(16)
            p.font.bold = True if i <= 2 else False
            p.font.color.rgb = self.COLORS['black']
            p.space_before = Pt(20)
            p.line_spacing = 1.2
            
    def _add_mckinsey_header(self, slide, title: str, subtitle: str):
        """Standardized consulting block header."""
        # Top line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2), self.prs.slide_width, Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS['light_blue']
        line.line.fill.background()

        # Headline
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.7))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['dark_blue']

        # Subtitle
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.5))
            p = sub.text_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(14)
            p.font.color.rgb = self.COLORS['gray']
            
            # Bottom line separator
            line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), self.prs.slide_width - Inches(1.0), Inches(0.02))
            line2.fill.solid()
            line2.fill.fore_color.rgb = self.COLORS['light_gray']
            line2.line.fill.background()

def generate_pptx_report(report_data: Dict[str, Any], charts: Dict[str, Any] = None) -> bytes:
    engine = PPTXStorytellingEngine()
    return engine.generate_report(report_data, charts)
